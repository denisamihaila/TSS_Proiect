"""
Script pentru generarea prezentarii TSS T1 - FitnessClassBooking
Ruleaza: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy
from pathlib import Path

# ── Paleta de culori ──────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0F, 0x17, 0x2A)   # fundal principal
BG_CARD     = RGBColor(0x1E, 0x29, 0x3B)   # card-uri / celule tabel
BG_HEADER   = RGBColor(0x0E, 0x1F, 0x38)   # header tabel
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)   # albastru deschis
ACCENT_TEAL = RGBColor(0x2D, 0xD4, 0xBF)   # teal
ACCENT_GRN  = RGBColor(0x34, 0xD3, 0x99)   # verde succes
ACCENT_AMB  = RGBColor(0xFB, 0xBF, 0x24)   # galben/amber
ACCENT_RED  = RGBColor(0xF8, 0x71, 0x71)   # rosu warning
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_L      = RGBColor(0xCB, 0xD5, 0xE1)   # text secundar deschis
GRAY_M      = RGBColor(0x64, 0x74, 0x8B)   # text tertiar
SLIDE_W     = Inches(13.33)
SLIDE_H     = Inches(7.5)


# ── Utilitare de baza ─────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """Seteaza culoarea de fundal a unui slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color: RGBColor, alpha=None):
    """Adauga un dreptunghi colorat pe slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Adauga un textbox simplu."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    """Adauga un paragraf intr-un text frame existent."""
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def add_multiline_box(slide, lines, x, y, w, h,
                      font_size=14, color=WHITE, bold_first=False):
    """Adauga un textbox cu mai multe linii."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = (bold_first and i == 0)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def add_table(slide, headers, rows, x, y, w, h,
              col_widths=None, header_color=BG_HEADER,
              alt_color=BG_CARD, hdr_font=13, row_font=12):
    """Adauga un tabel stilizat."""
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl = slide.shapes.add_table(nrows, ncols,
                                  Inches(x), Inches(y),
                                  Inches(w), Inches(h)).table
    # Latimi coloane
    if col_widths:
        total_inches = w
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    def style_cell(cell, text, bg, fg=WHITE, bold=False, sz=13, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg
        run.font.name = "Calibri"

    # Header
    for j, h_text in enumerate(headers):
        style_cell(tbl.cell(0, j), h_text, header_color,
                   fg=ACCENT_BLUE, bold=True, sz=hdr_font)

    # Randuri
    for i, row in enumerate(rows):
        bg = BG_CARD if i % 2 == 0 else RGBColor(0x27, 0x35, 0x49)
        for j, cell_text in enumerate(row):
            fg = WHITE
            # Colorare speciala pentru coloane de rezultate
            if isinstance(cell_text, tuple):
                cell_text, fg = cell_text
            style_cell(tbl.cell(i + 1, j), str(cell_text), bg,
                       fg=fg, sz=row_font, align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    return tbl


def accent_bar(slide, color=ACCENT_BLUE, y=0.85, height=0.04):
    """Linie de accent colorata orizontala."""
    add_rect(slide, 0, y, 13.33, height, color)


def slide_number(slide, num):
    """Adauga numarul slide-ului in coltul din dreapta jos."""
    add_textbox(slide, str(num), 12.5, 7.1, 0.6, 0.3,
                font_size=10, color=GRAY_M, align=PP_ALIGN.RIGHT)


def footer_line(slide, text="TSS T1 – FitnessClassBooking", color=GRAY_M):
    """Adauga un footer discret."""
    add_rect(slide, 0, 7.2, 13.33, 0.003, GRAY_M)
    add_textbox(slide, text, 0.3, 7.22, 12.0, 0.25,
                font_size=9, color=GRAY_M, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE-URI
# ═══════════════════════════════════════════════════════════════════════════════

def slide_1_title(prs):
    """Slide 1 – Titlu."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG_DARK)

    # Gradient decorativ stanga
    add_rect(slide, 0, 0, 0.5, 7.5, ACCENT_BLUE)
    add_rect(slide, 0.5, 0, 0.06, 7.5, ACCENT_TEAL)

    # Bloc decorativ dreapta
    add_rect(slide, 10.5, 0, 2.83, 7.5, BG_CARD)
    add_rect(slide, 10.44, 0, 0.06, 7.5, ACCENT_BLUE)

    # Badge disciplina
    add_rect(slide, 1.0, 1.2, 3.5, 0.5, ACCENT_BLUE)
    add_textbox(slide, "Testarea Sistemelor Software", 1.05, 1.25, 3.4, 0.4,
                font_size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    # Titlu principal
    add_textbox(slide, "TSS T1", 1.0, 1.95, 9.0, 1.0,
                font_size=52, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)

    # Subtitlu
    add_textbox(slide, "Testare unitară în Python", 1.0, 2.95, 9.2, 0.8,
                font_size=32, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    # Linie separator
    add_rect(slide, 1.0, 3.85, 8.0, 0.03, ACCENT_TEAL)

    # Detalii proiect
    details = [
        "Proiect:   FitnessClassBooking",
        "Metodă:   evaluate_client_package(session_history, package_sessions, has_membership)",
        "Framework:  pytest 9.0.3  |  coverage.py 7.13.5  |  mutmut 2.5.1  |  Cosmic Ray",
    ]
    for i, line in enumerate(details):
        add_textbox(slide, line, 1.0, 4.05 + i * 0.52, 9.5, 0.48,
                    font_size=15, color=GRAY_L if i > 0 else WHITE,
                    bold=(i == 0))

    # Metrici card-uri dreapta
    metrics = [
        ("100", "teste principale"),
        ("100%", "coverage"),
        ("0", "mutanți survived"),
    ]
    for i, (val, lbl) in enumerate(metrics):
        add_rect(slide, 10.7, 1.5 + i * 1.7, 2.3, 1.4, BG_DARK)
        add_textbox(slide, val, 10.7, 1.55 + i * 1.7, 2.3, 0.75,
                    font_size=34, bold=True, color=ACCENT_GRN, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, 10.7, 2.3 + i * 1.7, 2.3, 0.4,
                    font_size=12, color=GRAY_L, align=PP_ALIGN.CENTER)

    footer_line(slide)
    slide_number(slide, 1)


def slide_2_requirements(prs):
    """Slide 2 – Cerinte T1 si mapare pe implementare."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    accent_bar(slide, ACCENT_TEAL)

    add_textbox(slide, "Cerințe T1 și mapare pe implementare",
                0.4, 0.1, 12.5, 0.7, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "Framework: pytest  |  Toate strategiile cerute sunt acoperite",
                0.4, 0.75, 12.0, 0.4, font_size=13, color=ACCENT_TEAL)

    # Tabel strategii
    add_textbox(slide, "Strategii de testare", 0.4, 1.25, 6.2, 0.35,
                font_size=14, bold=True, color=ACCENT_BLUE)

    strat_headers = ["Strategie", "Fișier"]
    strat_rows = [
        ("Clase de echivalență", "test_equivalence_partitioning.py"),
        ("Valori de frontieră", "test_boundary_value_analysis.py"),
        ("Coverage instrucțiune/decizie/condiție", "test_coverage.py"),
        ("Circuite independente", "test_independent_circuits.py"),
        ("Mutation testing", "test_mutation.py"),
    ]
    add_table(slide, strat_headers, strat_rows,
              0.4, 1.65, 6.3, 2.9,
              col_widths=[3.5, 2.8],
              hdr_font=12, row_font=11)

    # Tabel cerinte structurale
    add_textbox(slide, "Cerințe structurale ale metodei", 7.0, 1.25, 6.0, 0.35,
                font_size=14, bold=True, color=ACCENT_BLUE)

    struct_headers = ["Cerință", "Implementare"]
    struct_rows = [
        ("≥ 3 parametri", "session_history, package_sessions, has_membership"),
        ("Instrucțiune repetitivă", "for session_status in session_history"),
        ("if cu else", "if session_status == 'attended': … else: …"),
        ("if fără else", "if has_membership: …"),
        ("Condiție simplă", "if has_membership"),
        ("Condiție compusă", "remaining_sessions == 0 and no_show == 0"),
    ]
    add_table(slide, struct_headers, struct_rows,
              7.0, 1.65, 6.0, 3.3,
              col_widths=[2.5, 3.5],
              hdr_font=12, row_font=10)

    # Nota Bool
    add_rect(slide, 0.4, 4.7, 12.5, 0.65, BG_CARD)
    add_rect(slide, 0.4, 4.7, 0.04, 0.65, ACCENT_AMB)
    add_textbox(slide, "Nota: Python – bool este subclasă a int → validarea separată a tipului bool este obligatorie "
                "(True/False ar putea trece ca 1/0 fără verificare explicită).",
                0.55, 4.74, 12.2, 0.55,
                font_size=12, color=ACCENT_AMB, italic=True)

    footer_line(slide)
    slide_number(slide, 2)


def slide_3_functionality(prs):
    """Slide 3 – Functionalitatea testata."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    accent_bar(slide, ACCENT_BLUE)

    add_textbox(slide, "Funcționalitatea testată – FitnessClassBooking",
                0.4, 0.1, 12.5, 0.7, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "Clasa modelează evaluarea unui pachet de ședințe de fitness cumpărat de un client",
                0.4, 0.75, 12.0, 0.4, font_size=13, color=ACCENT_BLUE)

    # Semnatura metodei - cod box
    add_rect(slide, 0.4, 1.25, 12.5, 1.05, RGBColor(0x0D, 0x14, 0x24))
    add_rect(slide, 0.4, 1.25, 0.04, 1.05, ACCENT_TEAL)
    code_lines = [
        "booking = FitnessClassBooking(class_name, instructor, price_per_session)",
        "result  = booking.evaluate_client_package(session_history, package_sessions, has_membership)",
    ]
    for i, line in enumerate(code_lines):
        add_textbox(slide, line, 0.55, 1.32 + i * 0.45, 12.2, 0.42,
                    font_size=13, color=ACCENT_TEAL)

    # Tabel statusuri
    add_textbox(slide, "Statusuri acceptate în session_history", 0.4, 2.45, 6.0, 0.35,
                font_size=14, bold=True, color=ACCENT_BLUE)
    status_headers = ["Status", "Semnificație", "Consumă ședință?"]
    status_rows = [
        ("attended", "Clientul a participat", ("Da", ACCENT_GRN)),
        ("no_show", "Absent fără anunț prealabil", ("Da", ACCENT_RED)),
        ("cancelled", "Anulat la timp", ("Nu", ACCENT_AMB)),
    ]
    add_table(slide, status_headers, status_rows,
              0.4, 2.85, 6.0, 1.8,
              col_widths=[1.8, 2.8, 1.4],
              hdr_font=12, row_font=12)

    # Exemplu concret
    add_textbox(slide, "Exemplu de calcul", 7.0, 2.45, 6.0, 0.35,
                font_size=14, bold=True, color=ACCENT_BLUE)
    add_rect(slide, 7.0, 2.85, 5.9, 1.8, RGBColor(0x0D, 0x14, 0x24))
    add_rect(slide, 7.0, 2.85, 0.04, 1.8, ACCENT_BLUE)
    ex_lines = [
        'Input: ["attended","attended","cancelled","no_show"],',
        '       package_sessions=5, has_membership=True',
        "─────────────────────────────────────────────",
        "used_sessions=3  remaining=2  cancelled=1",
        "cost = 5×50×0.80 = 200.0  |  status = active",
    ]
    for i, line in enumerate(ex_lines):
        color = ACCENT_AMB if i == 4 else (GRAY_M if i == 2 else GRAY_L)
        add_textbox(slide, line, 7.1, 2.92 + i * 0.32, 5.7, 0.3,
                    font_size=11, color=color)

    # Statusuri finale
    add_textbox(slide, "Status final al pachetului", 0.4, 4.75, 12.5, 0.35,
                font_size=14, bold=True, color=ACCENT_BLUE)
    status_cards = [
        ("active", "Mai are ședințe disponibile", ACCENT_GRN),
        ("completed_successfully", "Terminat fără no_show", ACCENT_BLUE),
        ("completed_with_absences", "Terminat cu cel puțin un no_show", ACCENT_RED),
    ]
    for i, (name, desc, color) in enumerate(status_cards):
        x = 0.4 + i * 4.3
        add_rect(slide, x, 5.15, 4.1, 0.9, BG_CARD)
        add_rect(slide, x, 5.15, 4.1, 0.08, color)
        add_textbox(slide, name, x + 0.1, 5.28, 3.9, 0.3,
                    font_size=12, bold=True, color=color)
        add_textbox(slide, desc, x + 0.1, 5.58, 3.9, 0.4,
                    font_size=11, color=GRAY_L)

    footer_line(slide)
    slide_number(slide, 3)


def slide_4_validations(prs):
    """Slide 4 – Validari si reguli de business."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    accent_bar(slide, ACCENT_AMB)

    add_textbox(slide, "Validări și reguli de business",
                0.4, 0.1, 12.5, 0.7, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "Orice intrare invalida aruncă ValueError – testele verifică mesajele publice de eroare",
                0.4, 0.75, 12.0, 0.4, font_size=13, color=ACCENT_AMB)

    # Constructor
    add_textbox(slide, "__init__(class_name, instructor, price_per_session)", 0.4, 1.3, 6.3, 0.4,
                font_size=14, bold=True, color=ACCENT_BLUE)
    ctor_headers = ["Parametru", "Regulă de validare"]
    ctor_rows = [
        ("class_name", "Trebuie să fie în {dance, pilates, yoga, zumba}"),
        ("instructor", "Șir nevid după strip()"),
        ("price_per_session", "Număr pozitiv; bool respins explicit"),
    ]
    add_table(slide, ctor_headers, ctor_rows,
              0.4, 1.75, 6.3, 1.8,
              col_widths=[2.2, 4.1],
              hdr_font=12, row_font=11)

    # Metoda
    add_textbox(slide, "evaluate_client_package(session_history, package_sessions, has_membership)",
                7.0, 1.3, 6.0, 0.4, font_size=14, bold=True, color=ACCENT_BLUE)
    meth_headers = ["Parametru", "Regulă de validare"]
    meth_rows = [
        ("session_history", "Trebuie să fie de tip list"),
        ("package_sessions", "int în [1, 20]; bool respins explicit"),
        ("has_membership", "Strict bool (nu int)"),
        ("statusuri în istoric", "Doar 'attended', 'no_show', 'cancelled'"),
        ("consum vs pachet", "used_sessions ≤ package_sessions"),
    ]
    add_table(slide, meth_headers, meth_rows,
              7.0, 1.75, 6.1, 2.2,
              col_widths=[2.3, 3.8],
              hdr_font=12, row_font=11)

    # Bool special case
    add_rect(slide, 0.4, 3.65, 8.0, 1.45, BG_CARD)
    add_rect(slide, 0.4, 3.65, 0.04, 1.45, ACCENT_AMB)
    add_textbox(slide, "Caz special Python: bool este subclasă de int",
                0.55, 3.7, 7.8, 0.4, font_size=13, bold=True, color=ACCENT_AMB)
    bool_cases = [
        "price_per_session=False  →  ValueError  (ar trece ca 0 fără validare)",
        "package_sessions=True   →  ValueError  (ar trece ca 1 fără validare)",
        "package_sessions=False  →  ValueError  (ar trece ca 0 fără validare)",
        "has_membership=1        →  ValueError  (nu este strict bool)",
    ]
    for i, case in enumerate(bool_cases):
        add_textbox(slide, case, 0.6, 4.12 + i * 0.25, 7.6, 0.24,
                    font_size=11, color=GRAY_L)

    # Frontiere package_sessions
    add_rect(slide, 8.7, 3.65, 4.5, 1.45, BG_CARD)
    add_rect(slide, 8.7, 3.65, 4.5, 0.04, ACCENT_RED)
    add_textbox(slide, "Frontiere package_sessions",
                8.8, 3.7, 4.2, 0.35, font_size=13, bold=True, color=ACCENT_RED)
    add_textbox(slide, "Invalid: 0  |  Valid: 1  ──  20  |  Invalid: 21",
                8.8, 4.1, 4.2, 0.35, font_size=13, color=GRAY_L, align=PP_ALIGN.CENTER)

    # Diagrama frontiere vizuala
    boundaries = [("0", ACCENT_RED), ("1", ACCENT_GRN), ("2", ACCENT_GRN),
                  ("···", GRAY_M), ("19", ACCENT_GRN), ("20", ACCENT_GRN), ("21", ACCENT_RED)]
    bx = 8.75
    for val, color in boundaries:
        w = 0.38 if val != "···" else 0.42
        add_rect(slide, bx, 4.55, w, 0.42, BG_DARK)
        add_rect(slide, bx, 4.55, w, 0.06, color)
        add_textbox(slide, val, bx, 4.62, w, 0.3,
                    font_size=11, color=color, align=PP_ALIGN.CENTER)
        bx += w + 0.07

    # Membership discount
    add_rect(slide, 0.4, 5.2, 5.8, 0.85, RGBColor(0x0D, 0x14, 0x24))
    add_rect(slide, 0.4, 5.2, 0.04, 0.85, ACCENT_GRN)
    add_textbox(slide, "Discount membership: 20%",
                0.55, 5.25, 5.6, 0.35, font_size=13, bold=True, color=ACCENT_GRN)
    add_textbox(slide, "total_cost = package_sessions × price_per_session × (1 − 0.20)",
                0.55, 5.6, 5.6, 0.35, font_size=12, color=GRAY_L)

    # Consum ≤ pachet
    add_rect(slide, 6.5, 5.2, 6.5, 0.85, RGBColor(0x0D, 0x14, 0x24))
    add_rect(slide, 6.5, 5.2, 0.04, 0.85, ACCENT_BLUE)
    add_textbox(slide, "Restricție: used_sessions ≤ package_sessions",
                6.65, 5.25, 6.2, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)
    add_textbox(slide, "Altfel → ValueError('invalid package data')",
                6.65, 5.6, 6.2, 0.35, font_size=12, color=GRAY_L)

    footer_line(slide)
    slide_number(slide, 4)


def slide_5_test_suite(prs):
    """Slide 5 – Proiectarea suitei principale de teste."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    accent_bar(slide, ACCENT_GRN)

    add_textbox(slide, "Proiectarea suitei principale de teste",
                0.4, 0.1, 12.5, 0.7, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "100 teste organizate pe tehnicile cerute în tema T1",
                0.4, 0.75, 12.0, 0.4, font_size=13, color=ACCENT_GRN)

    # Tabel principal
    tbl_headers = ["Fișier", "Strategie", "Nr. teste", "Rol principal"]
    tbl_rows = [
        ("test_equivalence_partitioning.py", "Clase de echivalență",
         ("23", ACCENT_BLUE), "Domenii valide/invalide pentru constructor și metodă"),
        ("test_boundary_value_analysis.py", "Valori de frontieră",
         ("16", ACCENT_BLUE), "Limite la 1, 20, prețuri minime și depășiri"),
        ("test_coverage.py", "Instrucțiune / Decizie / Condiție",
         ("40", ACCENT_BLUE), "Execuția tuturor ramurilor relevante"),
        ("test_independent_circuits.py", "Circuite independente",
         ("10", ACCENT_BLUE), "Drumuri reprezentative prin fluxul metodei"),
        ("test_mutation.py", "Teste orientate pe mutanți",
         ("11", ACCENT_BLUE), "Întărirea suitei pentru mutation testing"),
    ]
    add_table(slide, tbl_headers, tbl_rows,
              0.4, 1.2, 12.5, 3.4,
              col_widths=[3.6, 2.7, 1.1, 5.1],
              hdr_font=13, row_font=11)

    # Total badge
    add_rect(slide, 0.4, 4.7, 12.5, 0.72, BG_CARD)
    add_rect(slide, 0.4, 4.7, 12.5, 0.05, ACCENT_GRN)
    add_textbox(slide, "TOTAL SUITE PRINCIPALĂ:", 0.7, 4.78, 4.0, 0.5,
                font_size=16, bold=True, color=WHITE)
    add_textbox(slide, "100 teste", 4.7, 4.73, 2.5, 0.6,
                font_size=24, bold=True, color=ACCENT_GRN, align=PP_ALIGN.CENTER)
    add_textbox(slide, "toate structurate explicit pe tehnicile studiate la curs și laborator",
                7.2, 4.78, 5.5, 0.5, font_size=13, color=GRAY_L)

    # Card-uri sumare per strategie
    strategy_cards = [
        ("Funcțional", "23 + 16 = 39", "Echivalență + Frontiere", ACCENT_BLUE),
        ("Structural", "40 + 10 = 50", "Coverage + Circuite", ACCENT_TEAL),
        ("Mutation", "11", "Orientate pe mutanți", ACCENT_AMB),
    ]
    for i, (label, count, desc, color) in enumerate(strategy_cards):
        x = 0.4 + i * 4.2
        add_rect(slide, x, 5.55, 4.0, 1.5, BG_CARD)
        add_rect(slide, x, 5.55, 4.0, 0.05, color)
        add_textbox(slide, label, x + 0.15, 5.65, 3.7, 0.4,
                    font_size=13, bold=True, color=color)
        add_textbox(slide, count, x + 0.15, 6.05, 3.7, 0.55,
                    font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc, x + 0.15, 6.6, 3.7, 0.35,
                    font_size=11, color=GRAY_L, align=PP_ALIGN.CENTER)

    footer_line(slide)
    slide_number(slide, 5)


def slide_6_functional(prs):
    """Slide 6 – Testare functionala: echivalenta si frontiere."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    accent_bar(slide, ACCENT_BLUE)

    add_textbox(slide, "Testare funcțională: Echivalență și Frontiere",
                0.4, 0.1, 12.5, 0.7, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "23 teste pentru echivalență  +  16 teste pentru valori de frontieră",
                0.4, 0.75, 12.0, 0.4, font_size=13, color=ACCENT_BLUE)

    # Clase valide
    add_textbox(slide, "Clase valide", 0.4, 1.25, 5.8, 0.38,
                font_size=14, bold=True, color=ACCENT_GRN)
    valid_classes = [
        "class_name ∈ {dance, pilates, yoga, zumba}",
        "Instructor: șir nevid (ex: 'Ana Pop')",
        "Preț: orice număr pozitiv (ex: 50.0)",
        "session_history: listă goală sau cu statusuri valide",
        "package_sessions: int în intervalul [1, 20]",
        "has_membership: strict True sau False",
        "Pachet activ / finalizat cu succes / cu absențe",
    ]
    for i, line in enumerate(valid_classes):
        add_rect(slide, 0.4, 1.68 + i * 0.38, 0.15, 0.25, ACCENT_GRN)
        add_textbox(slide, line, 0.65, 1.68 + i * 0.38, 5.4, 0.34,
                    font_size=11, color=GRAY_L)

    # Clase invalide
    add_textbox(slide, "Clase invalide", 6.5, 1.25, 6.0, 0.38,
                font_size=14, bold=True, color=ACCENT_RED)
    invalid_classes = [
        "class_name: 'spinning', None, 123",
        "Instructor: '', '   ', 42",
        "Preț: 0, -10, True, False, 'cincizeci'",
        "session_history: None, tuple, string",
        "package_sessions: 0, 21, True, False, 5.0",
        "has_membership: 1, 0, 'yes', None",
        "Status: 'present', 'absent' (necunoscut)",
        "Consum > pachet: [attended×6] cu package=5",
    ]
    for i, line in enumerate(invalid_classes):
        add_rect(slide, 6.5, 1.68 + i * 0.38, 0.15, 0.25, ACCENT_RED)
        add_textbox(slide, line, 6.75, 1.68 + i * 0.38, 5.6, 0.34,
                    font_size=11, color=GRAY_L)

    # Separator vertical
    add_rect(slide, 6.25, 1.25, 0.04, 3.3, GRAY_M)

    # Frontiere - vizual
    add_textbox(slide, "Frontierele package_sessions", 0.4, 4.7, 8.0, 0.4,
                font_size=14, bold=True, color=ACCENT_AMB)

    bvals = [
        ("0", ACCENT_RED, "Invalid"),
        ("1", ACCENT_GRN, "Min valid"),
        ("2", ACCENT_GRN, "Min+1"),
        ("·  ·  ·", GRAY_M, ""),
        ("19", ACCENT_GRN, "Max−1"),
        ("20", ACCENT_GRN, "Max valid"),
        ("21", ACCENT_RED, "Invalid"),
    ]
    bx = 0.4
    for val, color, lbl in bvals:
        w = 1.4 if val == "·  ·  ·" else 1.5
        add_rect(slide, bx, 5.15, w - 0.1, 0.65, BG_CARD)
        add_rect(slide, bx, 5.15, w - 0.1, 0.08, color)
        add_textbox(slide, val, bx, 5.25, w - 0.1, 0.3,
                    font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        if lbl:
            add_textbox(slide, lbl, bx, 5.55, w - 0.1, 0.2,
                        font_size=9, color=GRAY_M, align=PP_ALIGN.CENTER)
        bx += w + 0.08

    # Frontiere pret
    add_textbox(slide, "Frontierele price_per_session", 0.4, 5.95, 6.0, 0.4,
                font_size=14, bold=True, color=ACCENT_AMB)
    price_cases = [
        ("< 0  →  ValueError", ACCENT_RED),
        ("= 0  →  ValueError", ACCENT_RED),
        ("> 0  →  Valid", ACCENT_GRN),
        ("bool  →  ValueError", ACCENT_RED),
    ]
    for i, (case, color) in enumerate(price_cases):
        add_rect(slide, 0.4 + i * 3.05, 6.38, 2.9, 0.5, BG_CARD)
        add_textbox(slide, case, 0.5 + i * 3.05, 6.45, 2.8, 0.38,
                    font_size=12, color=color, align=PP_ALIGN.CENTER)

    footer_line(slide)
    slide_number(slide, 6)


def slide_7_structural(prs):
    """Slide 7 – Testare structurala, CFG si coverage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    accent_bar(slide, ACCENT_TEAL)

    add_textbox(slide, "Testare structurală, CFG și Coverage",
                0.4, 0.1, 12.5, 0.7, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "40 teste de coverage  +  10 circuite independente  →  100% statement & branch coverage",
                0.4, 0.75, 12.5, 0.4, font_size=13, color=ACCENT_TEAL)

    # Ramuri acoperite
    add_textbox(slide, "Ramuri testate explicit", 0.4, 1.25, 6.0, 0.38,
                font_size=14, bold=True, color=ACCENT_BLUE)
    branches = [
        ("Validări constructor", "class_name, instructor, price invalide"),
        ("Validări metodă", "session_history, package_sessions, has_membership invalide"),
        ("Loop: status attended", "incrementează attended și used_sessions"),
        ("Loop: status no_show", "incrementează no_show și used_sessions"),
        ("Loop: status cancelled", "incrementează cancelled, nu consumă ședință"),
        ("Status invalid în buclă", "ValueError la status necunoscut"),
        ("Depășire pachet", "used_sessions > package_sessions → ValueError"),
        ("With membership", "total_cost × 0.80"),
        ("Without membership", "total_cost fără discount"),
        ("Status: active", "remaining_sessions > 0"),
        ("Status: completed_successfully", "remaining=0 AND no_show=0"),
        ("Status: completed_with_absences", "remaining=0 AND no_show>0"),
    ]
    for i, (branch, desc) in enumerate(branches):
        y = 1.68 + i * 0.42
        add_rect(slide, 0.4, y + 0.08, 0.12, 0.18, ACCENT_TEAL)
        add_textbox(slide, branch, 0.62, y, 2.4, 0.4,
                    font_size=10, bold=True, color=WHITE)
        add_textbox(slide, desc, 3.05, y, 3.3, 0.4,
                    font_size=10, color=GRAY_L)

    # Coverage result
    add_textbox(slide, "Rezultat Coverage", 6.7, 1.25, 6.4, 0.38,
                font_size=14, bold=True, color=ACCENT_BLUE)

    cov_headers = ["Metrica", "Valoare", "Status"]
    cov_rows = [
        ("Statements", "43 / 43", ("100%", ACCENT_GRN)),
        ("Missing", "0", ("", ACCENT_GRN)),
        ("Branches", "26 / 26", ("100%", ACCENT_GRN)),
        ("Partial branches", "0", ("", ACCENT_GRN)),
        ("Overall coverage", "fitness_class_booking.py", ("100%", ACCENT_GRN)),
    ]
    add_table(slide, cov_headers, cov_rows,
              6.7, 1.68, 6.4, 2.5,
              col_widths=[2.8, 2.0, 1.6],
              hdr_font=12, row_font=12)

    # Nota coverage
    add_rect(slide, 6.7, 4.28, 6.4, 0.65, BG_CARD)
    add_rect(slide, 6.7, 4.28, 0.04, 0.65, ACCENT_AMB)
    add_textbox(slide, "Coverage 100% arată că suita execută toate zonele "
                "din cod, dar nu garantează singur că testele sunt complete "
                "→ mutation testing completează imaginea.",
                6.85, 4.33, 6.1, 0.55,
                font_size=11, color=GRAY_L, italic=True)

    # CFG si Cause-Effect Graph
    add_textbox(slide, "Diagrame realizate cu tool dedicat (Draw.io)",
                6.7, 5.05, 6.4, 0.38,
                font_size=14, bold=True, color=ACCENT_BLUE)

    diag = [
        ("cfg_diagrama.drawio.png", "Control Flow Graph", "Fluxul complet al metodei evaluate_client_package"),
        ("cause_effect_graph.png", "Cause-Effect Graph", "Relații între condiții, reguli business și rezultate"),
    ]
    for i, (fname, title, desc) in enumerate(diag):
        y = 5.48 + i * 0.82
        add_rect(slide, 6.7, y, 6.4, 0.7, BG_CARD)
        add_rect(slide, 6.7, y, 0.04, 0.7, ACCENT_TEAL)
        add_textbox(slide, title, 6.85, y + 0.05, 6.1, 0.3,
                    font_size=12, bold=True, color=ACCENT_TEAL)
        add_textbox(slide, desc + f"  [{fname}]", 6.85, y + 0.35, 6.1, 0.28,
                    font_size=10, color=GRAY_L)

    footer_line(slide)
    slide_number(slide, 7)


def slide_8_mutation(prs):
    """Slide 8 – Mutation testing: mutmut si Cosmic Ray."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    accent_bar(slide, ACCENT_AMB)

    add_textbox(slide, "Mutation Testing: mutmut și Cosmic Ray",
                0.4, 0.1, 12.5, 0.7, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "Două instrumente complementare – operatori diferiți, perspectivă independentă",
                0.4, 0.75, 12.0, 0.4, font_size=13, color=ACCENT_AMB)

    # mutmut card
    add_rect(slide, 0.4, 1.25, 5.9, 4.6, BG_CARD)
    add_rect(slide, 0.4, 1.25, 5.9, 0.06, ACCENT_AMB)
    add_textbox(slide, "mutmut 2.5.1  (WSL Ubuntu 24.04.1)",
                0.6, 1.33, 5.5, 0.42, font_size=15, bold=True, color=ACCENT_AMB)

    mutmut_stats = [
        ("Mutanți verificați", "95 / 95", WHITE),
        ("Killed", "85", ACCENT_GRN),
        ("Suspicious", "10", ACCENT_AMB),
        ("Survived", "0", ACCENT_GRN),
        ("Timeout / Skipped", "0 / 0", GRAY_L),
    ]
    for i, (label, val, color) in enumerate(mutmut_stats):
        y = 1.85 + i * 0.65
        add_textbox(slide, label, 0.65, y, 3.2, 0.55,
                    font_size=13, color=GRAY_L)
        add_textbox(slide, val, 3.85, y, 2.1, 0.55,
                    font_size=18, bold=True, color=color, align=PP_ALIGN.RIGHT)

    add_rect(slide, 0.6, 5.05, 5.5, 0.04, GRAY_M)
    add_textbox(slide, "Suspicious ≠ Survived: indică rulări mai lente decât baseline, nu mutanți supraviețuitori.",
                0.6, 5.12, 5.5, 0.55, font_size=10, color=GRAY_M, italic=True)

    # Cosmic Ray card
    add_rect(slide, 6.8, 1.25, 6.1, 4.6, BG_CARD)
    add_rect(slide, 6.8, 1.25, 6.1, 0.06, ACCENT_BLUE)
    add_textbox(slide, "Cosmic Ray  (analiză suplimentară)",
                7.0, 1.33, 5.7, 0.42, font_size=15, bold=True, color=ACCENT_BLUE)

    cr_stats = [
        ("Mutanți generați", "166", WHITE),
        ("Killed", "157", ACCENT_GRN),
        ("Survived", "9", ACCENT_AMB),
        ("Rată supraviețuire", "5.42%", ACCENT_AMB),
        ("Scor kill", "94.58%", ACCENT_GRN),
    ]
    for i, (label, val, color) in enumerate(cr_stats):
        y = 1.85 + i * 0.65
        add_textbox(slide, label, 7.05, y, 3.5, 0.55,
                    font_size=13, color=GRAY_L)
        add_textbox(slide, val, 10.55, y, 2.1, 0.55,
                    font_size=18, bold=True, color=color, align=PP_ALIGN.RIGHT)

    add_rect(slide, 7.0, 5.05, 5.7, 0.04, GRAY_M)
    add_textbox(slide, "Cosmic Ray generează mai mulți mutanți cu operatori diferiți față de mutmut "
                "→ verificare complementară a robusteții suitei.",
                7.0, 5.12, 5.7, 0.55, font_size=10, color=GRAY_M, italic=True)

    # Test orientate pe mutanti
    add_rect(slide, 0.4, 6.0, 12.5, 1.0, RGBColor(0x0D, 0x14, 0x24))
    add_rect(slide, 0.4, 6.0, 0.04, 1.0, ACCENT_GRN)
    add_textbox(slide, "test_mutation.py (11 teste) – verifică explicit:",
                0.6, 6.05, 12.0, 0.38, font_size=13, bold=True, color=ACCENT_GRN)
    mut_points = ("costul se calculează pe pachetul complet, nu pe ședințele folosite  |  "
                  "discount de exact 20%  |  cancelled nu consumă ședință  |  "
                  "statusul final depinde de remaining_sessions ȘI no_show  |  mesaje de eroare stabile")
    add_textbox(slide, mut_points, 0.6, 6.45, 12.0, 0.45,
                font_size=11, color=GRAY_L)

    footer_line(slide)
    slide_number(slide, 8)


def slide_9_ai(prs):
    """Slide 9 – Utilizarea AI si comparatia suitei AI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    accent_bar(slide, ACCENT_TEAL)

    add_textbox(slide, "Utilizarea AI și suita asistată",
                0.4, 0.1, 12.5, 0.7, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "ChatGPT/Codex utilizat ca instrument de analiză și verificare – nu de înlocuire",
                0.4, 0.75, 12.0, 0.4, font_size=13, color=ACCENT_TEAL)

    # Cum a fost utilizat AI
    add_textbox(slide, "Activitati AI", 0.4, 1.25, 5.8, 0.38,
                font_size=14, bold=True, color=ACCENT_BLUE)
    ai_uses = [
        "Revizuirea suitei proprii de teste",
        "Identificarea cazurilor-limită (bool vs int)",
        "Propunerea testelor suplimentare pentru mutanți",
        "Construirea suitei independente în teste_ai/",
        "Compararea suitei proprii cu suita generată",
        "Formularea interpretărilor pentru coverage și mutation testing",
    ]
    for i, use in enumerate(ai_uses):
        add_rect(slide, 0.4, 1.7 + i * 0.45, 0.18, 0.25, ACCENT_TEAL)
        add_textbox(slide, use, 0.68, 1.7 + i * 0.45, 5.5, 0.38,
                    font_size=12, color=GRAY_L)

    # Suita AI fișiere
    add_textbox(slide, "Suita AI – teste_ai/  (70 teste)", 6.6, 1.25, 6.5, 0.38,
                font_size=14, bold=True, color=ACCENT_BLUE)
    ai_headers = ["Fișier", "Nr.", "Stil"]
    ai_rows = [
        ("test_ai_generated_booking.py", "13", "Scenarii business cu dataclass"),
        ("test_ai_equivalence_and_validation.py", "21", "Parametrizări echivalență+validare"),
        ("test_ai_boundary_and_structural.py", "22", "Limite, ramuri, rotunjire"),
        ("test_ai_paths_and_mutation_focus.py", "14", "Proprietăți, drumuri, mutații"),
    ]
    add_table(slide, ai_headers, ai_rows,
              6.6, 1.68, 6.5, 2.2,
              col_widths=[3.5, 0.6, 2.4],
              hdr_font=12, row_font=11)

    # Comparatie
    add_textbox(slide, "Comparație suita proprie vs suita AI", 0.4, 4.5, 12.5, 0.38,
                font_size=14, bold=True, color=ACCENT_BLUE)
    comp_headers = ["Criteriu", "Suita proprie", "Suita AI"]
    comp_rows = [
        ("Număr teste", "100", "70"),
        ("Organizare", "Pe tehnici de testare", "Pe scenarii, validări, proprietăți"),
        ("Stil", "Explicit, didactic", "Compact, parametrizat"),
        ("Scop", "Demonstrarea strategiilor cerute", "Perspectivă independentă"),
        ("Coverage", "100% (confirmat)", "Confirmă comportamentul"),
    ]
    add_table(slide, comp_headers, comp_rows,
              0.4, 4.93, 12.5, 2.2,
              col_widths=[2.5, 4.0, 6.0],
              hdr_font=12, row_font=11)

    # Studii de caz
    add_rect(slide, 6.6, 3.98, 6.5, 0.42, BG_CARD)
    add_rect(slide, 6.6, 3.98, 0.04, 0.42, ACCENT_AMB)
    add_textbox(slide, "Studii de caz AI: bool vs int (True/False ca 1/0)  |  "
                "egalitate string via ==, nu is  |  statusuri construite dinamic",
                6.75, 4.01, 6.2, 0.36,
                font_size=10, color=ACCENT_AMB)

    footer_line(slide)
    slide_number(slide, 9)


def slide_10_results(prs):
    """Slide 10 – Rezultate finale si livrabile."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    # Bara colorata sus
    add_rect(slide, 0, 0, 13.33, 0.06, ACCENT_GRN)

    add_textbox(slide, "Rezultate finale și livrabile",
                0.4, 0.15, 12.5, 0.7, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "Proiectul respectă cerința T1 – testare funcțională + structurală + mutation + AI",
                0.4, 0.8, 12.0, 0.4, font_size=13, color=ACCENT_GRN)

    # Metrici principale – 6 card-uri
    metrics = [
        ("100", "teste principale", ACCENT_GRN),
        ("70", "teste AI", ACCENT_BLUE),
        ("100%", "statement coverage", ACCENT_GRN),
        ("100%", "branch coverage", ACCENT_GRN),
        ("0", "mutanți survived\n(mutmut)", ACCENT_GRN),
        ("94.58%", "kill rate\n(Cosmic Ray)", ACCENT_AMB),
    ]
    cards_per_row = 3
    for i, (val, lbl, color) in enumerate(metrics):
        row = i // cards_per_row
        col = i % cards_per_row
        x = 0.4 + col * 4.28
        y = 1.35 + row * 2.0
        add_rect(slide, x, y, 4.0, 1.75, BG_CARD)
        add_rect(slide, x, y, 4.0, 0.07, color)
        add_textbox(slide, val, x + 0.1, y + 0.2, 3.8, 0.9,
                    font_size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, x + 0.1, y + 1.1, 3.8, 0.55,
                    font_size=12, color=GRAY_L, align=PP_ALIGN.CENTER)

    # Livrabile
    add_textbox(slide, "Livrabile", 0.4, 5.5, 5.8, 0.38,
                font_size=14, bold=True, color=ACCENT_BLUE)
    deliverables = [
        "README.md – documentație completă (713 linii)",
        "fitness_class_booking.py – clasa testată",
        "test_*.py (5 fișiere) – suita principală 100 teste",
        "teste_ai/ (4 fișiere) – suita AI 70 teste",
        "cfg_diagrama.drawio.png, cause_effect_graph.png",
        "screenshots/ (7) + logs/ + cosmic_ray/",
        "TSS_T1_FitnessClassBooking.pptx – prezentarea",
    ]
    for i, item in enumerate(deliverables):
        add_rect(slide, 0.4, 5.93 + i * 0.22, 0.12, 0.15, ACCENT_BLUE)
        add_textbox(slide, item, 0.62, 5.9 + i * 0.22, 5.5, 0.2,
                    font_size=10, color=GRAY_L)

    # Concluzie
    add_rect(slide, 6.5, 5.5, 6.5, 1.7, BG_CARD)
    add_rect(slide, 6.5, 5.5, 6.5, 0.06, ACCENT_GRN)
    add_textbox(slide, "Concluzie",
                6.65, 5.58, 6.2, 0.38, font_size=14, bold=True, color=ACCENT_GRN)
    concl = ("Metoda evaluate_client_package permite verificarea coerentă a claselor de echivalență, "
             "valorilor de frontieră, acoperirii codului, circuitelor independente și comportamentului "
             "în fața mutațiilor. Rezultatele obținute indică o suită stabilă și bine focalizată pe "
             "contractul public al clasei.")
    add_textbox(slide, concl, 6.65, 6.0, 6.2, 1.1,
                font_size=11, color=GRAY_L)

    footer_line(slide)
    slide_number(slide, 10)


# ── Configurare dimensiune slide 13.33 × 7.5 inch (widescreen 16:9) ──────────

def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Generez slide-urile...")
    slide_1_title(prs)
    print("  [1/10] Titlu")
    slide_2_requirements(prs)
    print("  [2/10] Cerinte")
    slide_3_functionality(prs)
    print("  [3/10] Functionalitate")
    slide_4_validations(prs)
    print("  [4/10] Validari")
    slide_5_test_suite(prs)
    print("  [5/10] Suita teste")
    slide_6_functional(prs)
    print("  [6/10] Testare functionala")
    slide_7_structural(prs)
    print("  [7/10] Testare structurala")
    slide_8_mutation(prs)
    print("  [8/10] Mutation testing")
    slide_9_ai(prs)
    print("  [9/10] AI")
    slide_10_results(prs)
    print("  [10/10] Rezultate finale")

    output = Path(__file__).with_name("TSS_T1_FitnessClassBooking.pptx")
    prs.save(output)
    print(f"\nPrezentare salvata: {output}")


if __name__ == "__main__":
    build_presentation()
